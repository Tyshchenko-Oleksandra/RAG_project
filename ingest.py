import os
from pathlib import Path
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation
from openai import OpenAI
import chromadb

load_dotenv()

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

DATA_DIR = Path("data")
CHROMA_DIR = "chroma_db"

chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)
collection = chroma_client.get_or_create_collection(name="university_notes")


def read_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    texts = []

    for page_number, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""

        if page_text.strip():
            texts.append(f"Page {page_number}:\n{page_text.strip()}")

    return "\n\n".join(texts)


def read_pptx(file_path: Path) -> str:
    presentation = Presentation(file_path)
    texts = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            texts.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))

    return "\n\n".join(texts)


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]

        if chunk.strip():
            chunks.append(chunk.strip())

        start = end - overlap

    return chunks


def create_embedding(text: str) -> list[float]:
    response = client.embeddings.create(
        model="text-embedding-3-small",
        input=text
    )
    return response.data[0].embedding


def ingest_file(file_path: Path, subject: str):
    print(f"\nReading: {file_path}")

    if file_path.suffix.lower() == ".pdf":
        text = read_pdf(file_path)
    elif file_path.suffix.lower() == ".pptx":
        text = read_pptx(file_path)
    else:
        print(f"Skipping unsupported file: {file_path}")
        return

    if not text.strip():
        print(f"No text found in {file_path.name}. Можливо, це скан або картинка.")
        return

    chunks = chunk_text(text)

    for index, chunk in enumerate(chunks):
        embedding = create_embedding(chunk)

        doc_id = f"{subject}-{file_path.stem}-{index}"

        collection.add(
            ids=[doc_id],
            embeddings=[embedding],
            documents=[chunk],
            metadatas=[{
                "subject": subject,
                "source": file_path.name,
                "chunk_index": index
            }]
        )

    print(f"Added {len(chunks)} chunks from {file_path.name}")


def ingest_all():
    if not DATA_DIR.exists():
        print("Папка data не знайдена.")
        return

    files_found = 0

    for subject_dir in DATA_DIR.iterdir():
        if subject_dir.is_dir():
            subject = subject_dir.name
            print(f"\nSubject: {subject}")

            files = list(subject_dir.rglob("*"))

            for file_path in files:
                if file_path.is_file():
                    print(f"Found file: {file_path}")

                    if file_path.suffix.lower() in [".pdf", ".pptx"]:
                        files_found += 1
                        ingest_file(file_path, subject)
                    else:
                        print(f"Skipping unsupported file type: {file_path.suffix}")

    print(f"\nTotal supported files found: {files_found}")
    print("Done!")

if __name__ == "__main__":
    ingest_all()
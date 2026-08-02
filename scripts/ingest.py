import os
from pathlib import Path
from dotenv import load_dotenv
from pypdf import PdfReader
from pptx import Presentation
from openai import OpenAI
import chromadb

load_dotenv()

client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

DATA_DIR = Path("data")
CHROMA_DIR = "chroma_db"

chroma_client = chromadb.PersistentClient(path=CHROMA_DIR)
collection = chroma_client.get_or_create_collection(name="university_notes")


def read_pdf(file_path: Path) -> list[dict]:
    reader = PdfReader(str(file_path))
    pages = []

    for page_number, page in enumerate(reader.pages, start=1):
        page_text = page.extract_text() or ""

        if page_text.strip():
            pages.append({
                "text": page_text.strip(),
                "page": page_number
            })

    return pages


def read_pptx(file_path: Path) -> list[dict]:
    presentation = Presentation(file_path)
    slides = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

        if slide_text:
            slides.append({
                "text": "\n".join(slide_text),
                "slide": slide_number
            })

    return slides


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 150) -> list[str]:
    chunks = []
    start = 0

    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]

        if chunk.strip():
            chunks.append(chunk.strip())

        start = end - overlap

    return chunks


def create_embedding(text: str) -> list[float]:
    response = client.embeddings.create(
        model="text-embedding-3-small",
        input=text
    )
    return response.data[0].embedding


def ingest_file(file_path: Path, subject: str):
    print(f"\nReading: {file_path}")

    if file_path.suffix.lower() == ".pdf":
        pages = read_pdf(file_path)

        if not pages:
            print(f"No text found in {file_path.name}. Можливо, це скан або картинка.")
            return

        for page_data in pages:
            page_number = page_data["page"]
            text = page_data["text"]

            chunks = chunk_text(text)

            for index, chunk in enumerate(chunks):
                embedding = create_embedding(chunk)

                doc_id = f"{subject}-{file_path.stem}-page-{page_number}-chunk-{index}"

                collection.upsert(
                    ids=[doc_id],
                    embeddings=[embedding],
                    documents=[chunk],
                    metadatas=[{
                        "subject": subject,
                        "source": file_path.name,
                        "page": page_number,
                        "chunk_index": index,
                        "content_type": "pdf"
                    }]
                )

        print(f"Added PDF pages from {file_path.name}")

    elif file_path.suffix.lower() == ".pptx":
        slides = read_pptx(file_path)

        if not slides:
            print(f"No text found in {file_path.name}.")
            return

        for slide_data in slides:
            slide_number = slide_data["slide"]
            text = slide_data["text"]

            chunks = chunk_text(text)

            for index, chunk in enumerate(chunks):
                embedding = create_embedding(chunk)

                doc_id = f"{subject}-{file_path.stem}-slide-{slide_number}-chunk-{index}"

                collection.upsert(
                    ids=[doc_id],
                    embeddings=[embedding],
                    documents=[chunk],
                    metadatas=[{
                        "subject": subject,
                        "source": file_path.name,
                        "slide": slide_number,
                        "chunk_index": index,
                        "content_type": "pptx"
                    }]
                )

        print(f"Added PPTX slides from {file_path.name}")

    else:
        print(f"Skipping unsupported file: {file_path}")

    print(f"Added {len(chunks)} chunks from {file_path.name}")


def ingest_all():
    if not DATA_DIR.exists():
        print("Папка data не знайдена.")
        return

    files_found = 0

    for subject_dir in DATA_DIR.iterdir():
        if subject_dir.is_dir():
            subject = subject_dir.name
            print(f"\nSubject: {subject}")

            files = list(subject_dir.rglob("*"))

            for file_path in files:
                if file_path.is_file():
                    print(f"Found file: {file_path}")

                    if file_path.suffix.lower() in [".pdf", ".pptx"]:
                        files_found += 1
                        ingest_file(file_path, subject)
                    else:
                        print(f"Skipping unsupported file type: {file_path.suffix}")

    print(f"\nTotal supported files found: {files_found}")
    print("Done!")

if __name__ == "__main__":
    ingest_all()